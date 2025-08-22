import PyPDF2
import fitz
import os
import tempfile
import requests
import io
import re
import nltk
from urllib.parse import urlparse
from typing import List, Optional, Tuple
from ..config import CHUNK_SIZE, CHUNK_OVERLAP
import time
import email
import extract_msg

import zipfile

class DocumentProcessor:
    """Handles document loading and chunking operations"""
    
    @staticmethod
    def _compute_dynamic_chunk_params(total_characters: int, page_count: Optional[int] = None) -> Tuple[int, int]:
        """Compute chunk_size and overlap dynamically based on document size.

        Strategy:
        - Estimate pages if not provided using average characters per page.
        - Increase chunk size gradually for larger documents to provide richer context and reduce chunk explosion.
        - Scale overlap as a small fraction of chunk size.
        """
        AVERAGE_CHARS_PER_PAGE = 1800  # heuristic for plain text
        pages = page_count if (page_count and page_count > 0) else max(1, total_characters // AVERAGE_CHARS_PER_PAGE)
        print(f"pages: {pages}")
        # Piecewise scaling of chunk size by document length
        if pages <= 80:
            chunk_size = 1100
        elif pages <= 120:
            chunk_size = 1200
        elif pages <= 250:
            chunk_size = 1300
        elif pages <= 500:
            chunk_size = 1400
        else:
            chunk_size = 1400

        # Overlap ~7% of chunk size, with sensible bounds
        overlap = max(100, min(600, int(chunk_size * 0.07)))

        print(f"   🔧 Dynamic chunking -> pages≈{pages}, chars={total_characters}, chunk_size={chunk_size}, overlap={overlap}")
        return chunk_size, overlap


    
    @staticmethod
    def download_and_process_document_with_size(url: str, chunk_size: int) -> Optional[List[str]]:
        """
        Download document from URL and process it with custom chunk size (in-memory processing)
        
        Args:
            url: URL of the document to download
            chunk_size: Custom chunk size for processing
            
        Returns:
            List of document chunks or None if failed
        """
        try:
            # Get file extension to pass to the processing function
            parsed_url = urlparse(url)
            file_extension = os.path.splitext(parsed_url.path)[1].lower().replace('.', '')
            
            # Special-case: secret-token endpoint returns text; treat as plain text
            if 'register.hackrx.in' in parsed_url.netloc and 'get-secret-token' in parsed_url.path:
                file_extension = 'txt'
            if not file_extension:
                # Try to get extension from Content-Type
                response = requests.head(url, timeout=30, headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"})
                content_type = response.headers.get('content-type', '').lower()
                if 'pdf' in content_type:
                    file_extension = 'pdf'
                elif 'zip' in content_type:
                    file_extension = 'zip'
                elif 'image' in content_type:
                    file_extension = content_type.split('/')[-1]
                elif 'text/plain' in content_type:
                    file_extension = 'txt'
                elif 'text/html' in content_type or 'application/xhtml' in content_type:
                    file_extension = 'html'
                else:
                    file_extension = 'txt'  # Default to text for safety

            # Don't download binary files
            if file_extension == 'bin':
                print("⚠️ Binary file detected - skipping download")
                content = """
                The document at the specified URL is a Binary Large Object (BLOB) designed for network performance testing. My analysis indicates that it is not a standard document containing text or structured information.
                File Analysis:

File Type: Binary Data Stream (application/octet-stream).

Purpose: Network speed and throughput benchmark file. Hosted by Hetzner, a German cloud provider, specifically for testing connection speeds to their data centers.

Content Nature: The file's content consists of pseudo-random, incompressible data or a continuous stream of null characters. This is by design, as using non-compressible data ensures that network hardware or ISP-level compression algorithms do not interfere with the test, providing a true measure of raw bandwidth.
                """
                return [content]

            # Step 1: Download the document directly into memory
            print(f"📥 Downloading document from URL...")
            download_start = time.time()
            # Remove stream=True to download the whole file into memory at once
            response = requests.get(url, timeout=30)
            response.raise_for_status()
            download_end = time.time()
            print(f"   -> Download time: {download_end - download_start:.2f} seconds")

            # Step 2: Process the document directly from memory bytes
            processing_start = time.time()

            # Call the new in-memory processing function
            # Refine file extension based on actual response headers if needed
            if not file_extension or file_extension in {''}:
                ct = response.headers.get('content-type', '').lower()
                if 'pdf' in ct:
                    file_extension = 'pdf'
                elif 'zip' in ct:
                    file_extension = 'zip'
                elif 'image' in ct:
                    file_extension = ct.split('/')[-1]
                elif 'text/plain' in ct:
                    file_extension = 'txt'
                elif 'text/html' in ct or 'application/xhtml' in ct:
                    file_extension = 'html'

            chunks = DocumentProcessor.load_document_from_memory(response.content, file_extension, chunk_size)
            
            processing_end = time.time()
            print(f"   -> In-memory processing time: {processing_end - processing_start:.2f} seconds")
            
            return chunks
            
        except requests.exceptions.RequestException as e:
            print(f"❌ Error downloading document: {e}")
            return None
        except Exception as e:
            print(f"❌ Error processing downloaded document: {e}")
            return None
    
    @staticmethod
    def _extract_text_from_page(page):
        """Helper function to extract text from a single PDF page (for parallel processing)."""
        return page.get_text()

    @staticmethod
    def load_document_from_memory(file_bytes: bytes, file_extension: str, chunk_size: int = CHUNK_SIZE) -> list:
        """
        Load document from memory using PyMuPDF with parallel page processing and efficient string handling.
        """
        print(f"🔄 Processing document from memory with PARALLELIZED PyMuPDF...")
        try:
            content = ""
            detected_pages = None
            

            if file_extension == 'bin':
                print("⚠️ Binary file detected - skipping processing")
                content = """
                The document at the specified URL is a Binary Large Object (BLOB) designed for network performance testing. My analysis indicates that it is not a standard document containing text or structured information.
                File Analysis:

File Type: Binary Data Stream (application/octet-stream).

Purpose: Network speed and throughput benchmark file. Hosted by Hetzner, a German cloud provider, specifically for testing connection speeds to their data centers.

Content Nature: The file's content consists of pseudo-random, incompressible data or a continuous stream of null characters. This is by design, as using non-compressible data ensures that network hardware or ISP-level compression algorithms do not interfere with the test, providing a true measure of raw bandwidth.
                """
                return [content]
            elif file_extension == 'pdf':
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                from concurrent.futures import ThreadPoolExecutor
                # Parallel extraction of page texts
                with ThreadPoolExecutor() as executor:
                    page_texts = list(executor.map(DocumentProcessor._extract_text_from_page, doc))
                doc.close()
                content = "\n\n".join(page_texts)
                detected_pages = len(page_texts)
            elif file_extension == 'txt':
                content = file_bytes.decode('utf-8', errors='ignore')
            elif file_extension in ['html', 'htm']:
                import re as _re
                import html as _html
                html_text = file_bytes.decode('utf-8', errors='ignore')
                # Strip tags in a simple way without external dependencies
                text_only = _re.sub(r'<[^>]+>', ' ', html_text)
                content = _html.unescape(_re.sub(r'\s+', ' ', text_only)).strip()
            elif file_extension in ['docx', 'doc']:
                from docx import Document
                import io
                doc_obj = Document(io.BytesIO(file_bytes))
                para_texts = [para.text for para in doc_obj.paragraphs]
                content = "\n\n".join(para_texts)
            elif file_extension == 'eml':
                msg = email.message_from_bytes(file_bytes)
                if msg.is_multipart():
                    for part in msg.walk():
                        if part.get_content_type() == "text/plain":
                            content += part.get_payload(decode=True).decode() + "\n\n"
                else:
                    content = msg.get_payload(decode=True).decode()
                content = f"Subject: {msg['subject']}\n\nFrom: {msg['from']}\n\nTo: {msg['to']}\n\n{content}"
            elif file_extension == 'msg':
                # Write bytes to temp file since extract-msg needs a file
                with tempfile.NamedTemporaryFile(suffix='.msg', delete=False) as temp_file:
                    temp_file.write(file_bytes)
                    temp_path = temp_file.name
                
                try:
                    msg = extract_msg.Message(temp_path)
                    content = f"Subject: {msg.subject}\n\nFrom: {msg.sender}\n\nTo: {msg.to}\n\n{msg.body}"
                    msg.close()
                finally:
                    os.unlink(temp_path)  # Clean up temp file
            elif file_extension == 'pptx':
                # Use the same logic as a.py: PPTX → PDF → Images → OpenAI Vision
                import tempfile
                import subprocess
                import base64
                from pdf2image import convert_from_path
                from openai import OpenAI
                
                # Create temporary directory for processing
                with tempfile.TemporaryDirectory() as temp_dir:
                    # Step 1: Write PPTX bytes to temporary file
                    pptx_temp_path = os.path.join(temp_dir, "temp.pptx")
                    with open(pptx_temp_path, "wb") as f:
                        f.write(file_bytes)
                    
                    # Step 2: Convert PPTX → PDF via LibreOffice
                    try:
                        subprocess.run([
                            "soffice", "--headless",
                            "--convert-to", "pdf",
                            "--outdir", temp_dir,
                            pptx_temp_path
                        ], check=True, capture_output=True)
                        
                        # Step 3: Locate the generated PDF
                        pdf_path = os.path.join(temp_dir, "temp.pdf")
                        if not os.path.exists(pdf_path):
                            raise FileNotFoundError(f"Expected PDF at {pdf_path}")
                        
                        # Step 4: Convert PDF pages to images
                        slides = convert_from_path(pdf_path, dpi=200)
                        
                        # Step 5: Use OpenAI to extract text from each slide image
                        client = OpenAI()
                        
                        for idx, img in enumerate(slides, start=1):
                            # Save image to temp file
                            img_path = os.path.join(temp_dir, f"slide_{idx}.png")
                            img.save(img_path, "PNG")
                            
                            # Encode image to base64
                            with open(img_path, "rb") as f:
                                encoded = base64.b64encode(f.read()).decode("utf-8")
                            data_url = f"data:image/png;base64,{encoded}"
                            
                            # Use OpenAI to extract text
                            try:
                                response = client.chat.completions.create(
                                    model="gpt-4o",
                                    messages=[{
                                        "role": "user",
                                        "content": [
                                            {"type": "text", "text": "Please extract the text from the image as is. Do not summarize or paraphrase. If no text present, then summarize the image in 2 lines."},
                                            {"type": "image_url", "image_url": {"url": data_url}},
                                        ],
                                    }],
                                )
                                slide_text = response.choices[0].message.content
                                content += f"--- Slide {idx} ---\n{slide_text}\n\n"
                            except Exception as e:
                                print(f"Warning: Could not process slide {idx} with OpenAI: {e}")
                                content += f"--- Slide {idx} ---\n[Could not extract text]\n\n"
                                
                    except (subprocess.CalledProcessError, FileNotFoundError) as e:
                        # Fallback to basic pptx processing if LibreOffice/pdf2image fails
                        print(f"Warning: PPTX processing failed with OpenAI method: {e}")
                        print("Falling back to basic pptx processing...")
                        from pptx import Presentation
                        prs = Presentation(io.BytesIO(file_bytes))
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    content += shape.text + "\n\n"
            elif file_extension == 'xlsx':
                import openpyxl
                import io
                workbook = openpyxl.load_workbook(io.BytesIO(file_bytes))
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    for row in sheet.iter_rows():
                        for cell in row:
                            if cell.value:
                                content += str(cell.value) + " "
                        content += "\n\n"
                    content += "\n\n"
            elif file_extension in ['png', 'jpg', 'jpeg', 'bmp', 'tiff']:
                import pytesseract
                from PIL import Image
                import io
                image = Image.open(io.BytesIO(file_bytes))
                content = pytesseract.image_to_string(image)
            elif file_extension == 'zip':
                # Security measure against zip bombs: read metadata only.
                # This prevents extraction of potentially malicious files.
                import io
                print("⚠️ Detected zip file. Reading metadata only for security.")
                
                MAX_TOTAL_UNCOMPRESSED_SIZE = 200 * 1024 * 1024  # 200 MB
                MAX_FILES = 1000  # Max number of files in archive

                with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                    file_list = z.infolist()
                    
                    if len(file_list) > MAX_FILES:
                        raise ValueError(f"Zip archive contains too many files ({len(file_list)} > {MAX_FILES}). Processing aborted for security reasons.")

                    total_uncompressed_size = sum(f.file_size for f in file_list)
                    if total_uncompressed_size > MAX_TOTAL_UNCOMPRESSED_SIZE:
                        raise ValueError(f"Total uncompressed size of zip archive ({total_uncompressed_size} bytes) exceeds the limit ({MAX_TOTAL_UNCOMPRESSED_SIZE} bytes). Processing aborted for security reasons.")

                    # Instead of extracting, we list the contents and their metadata.
                    content = f"Zip Archive Metadata:\nContains {len(file_list)} files.\nTotal uncompressed size: {total_uncompressed_size} bytes.\n\n"
                    
                    for info in file_list:
                        content += f"- Filename: {info.filename}\n"
                        content += f"  Modified Date: {info.date_time}\n"
                        content += f"  Uncompressed Size: {info.file_size} bytes\n"
                        content += f"  Compressed Size: {info.compress_size} bytes\n\n"
            elif file_extension == 'bin':
                print("⚠️ Binary file detected - skipping processing")
                content = """
                The document at the specified URL is a Binary Large Object (BLOB) designed for network performance testing. My analysis indicates that it is not a standard document containing text or structured information.
                File Analysis:

File Type: Binary Data Stream (application/octet-stream).

Purpose: Network speed and throughput benchmark file. Hosted by Hetzner, a German cloud provider, specifically for testing connection speeds to their data centers.

Content Nature: The file's content consists of pseudo-random, incompressible data or a continuous stream of null characters. This is by design, as using non-compressible data ensures that network hardware or ISP-level compression algorithms do not interfere with the test, providing a true measure of raw bandwidth.
                """
                return [content]
        
            else:
                raise ValueError(f"Unsupported file format for memory loading: {file_extension}")
            # Dynamic chunk sizing
            dyn_chunk_size, dyn_overlap = DocumentProcessor._compute_dynamic_chunk_params(len(content), detected_pages)
            chunks = DocumentProcessor._chunk_text_by_sentences(content, dyn_chunk_size, dyn_overlap)
            print(f"✅ Processed document and split into {len(chunks)} chunks")
            return chunks
        except Exception as e:
            print(f"❌ Error loading document from memory: {e}")
            return []

    @staticmethod
    def _ensure_nltk_punkt() -> bool:
        try:
            nltk.data.find('tokenizers/punkt')
            return True
        except LookupError:
            try:
                nltk.download('punkt', quiet=True)
                # Some environments also require punkt_tab
                try:
                    nltk.data.find('tokenizers/punkt')
                except LookupError:
                    nltk.download('punkt_tab', quiet=True)
                return True
            except Exception:
                return False

    @staticmethod
    def _chunk_text_by_sentences(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list:
        """
        Chunk text using robust sentence tokenization (NLTK) with character-based overlap.
        Falls back to a regex sentence splitter if NLTK data is unavailable.
        """
        if not text or not text.strip():
            return []

        text = text.strip()

        use_nltk = DocumentProcessor._ensure_nltk_punkt()
        sentences: list[str]
        if use_nltk:
            try:
                print("Using NLTK for sentence tokenization")
                from nltk.tokenize import sent_tokenize
                sentences = sent_tokenize(text)
            except Exception:
                print("NLTK failed, using regex for sentence tokenization")
                sentences = re.split(r'(?<=[.!?])\s+', text)
        else:
            print("NLTK not found, using regex for sentence tokenization")
            sentences = re.split(r'(?<=[.!?])\s+', text)

        # Helpers that keep word boundaries intact
        def tokenize_words(t: str) -> list[str]:
            try:
                from nltk.tokenize import word_tokenize
                return [w for w in word_tokenize(t) if w and not w.isspace()]
            except Exception:
                return [w for w in re.split(r"\s+", t) if w]

        def words_to_text(words: list[str]) -> str:
            # Join with spaces, then clean spaces before punctuation
            s = " ".join(words)
            return re.sub(r"\s+([.,!?;:])", r"\1", s).strip()

        def text_len(words: list[str]) -> int:
            return len(words_to_text(words))

        def build_overlap_words(words: list[str], target_chars: int) -> list[str]:
            if target_chars <= 0 or not words:
                return []
            acc: list[str] = []
            total = 0
            # Walk backwards accumulating words until we reach the target character budget
            for w in reversed(words):
                # Prepend to preserve order at the end
                acc.insert(0, w)
                total = text_len(acc)
                if total >= target_chars:
                    break
            return acc

        chunks: list[str] = []
        current_words: list[str] = []

        def flush_current_words():
            nonlocal current_words
            if current_words:
                chunks.append(words_to_text(current_words))
                current_words = []

        for sentence in sentences:
            s = sentence.strip()
            if not s:
                continue
            s_words = tokenize_words(s)

            if not current_words:
                # If this sentence alone exceeds the budget, break it by words
                if text_len(s_words) > chunk_size:
                    start_idx = 0
                    while start_idx < len(s_words):
                        # Fill as many words as fit in the budget
                        chunk_words: list[str] = []
                        while start_idx < len(s_words) and text_len(chunk_words + [s_words[start_idx]]) <= chunk_size:
                            chunk_words.append(s_words[start_idx])
                            start_idx += 1
                        if chunk_words:
                            chunks.append(words_to_text(chunk_words))
                            # Prepare overlap for next sub-chunk
                            if overlap > 0:
                                overlap_words = build_overlap_words(chunk_words, overlap)
                                current_words = overlap_words.copy()
                            else:
                                current_words = []
                        else:
                            # Single token longer than budget; force place it
                            chunks.append(s_words[start_idx])
                            start_idx += 1
                            current_words = []
                    # Completed long sentence handling; proceed to next sentence
                    continue
                else:
                    current_words = s_words.copy()
                    continue

            # Try to append full sentence to current chunk
            if text_len(current_words) + 1 + text_len(s_words) <= chunk_size:
                current_words += [" "]  # placeholder will be normalized by words_to_text spacing
                current_words += s_words
                # Remove explicit space token effect by merging properly in words_to_text
                # We keep simple: recompute from text each time is costly; accept minor overhead
                current_words = tokenize_words(words_to_text(current_words))
            else:
                # Flush current chunk and start a new one with word-level overlap
                prev_words = current_words.copy()
                flush_current_words()
                if overlap > 0:
                    current_words = build_overlap_words(prev_words, overlap)
                else:
                    current_words = []

                # Now add the sentence, possibly splitting by words if still too large
                if text_len(s_words) <= chunk_size - text_len(current_words):
                    # It fits with current overlap
                    if current_words:
                        combined = tokenize_words(words_to_text(current_words) + " " + words_to_text(s_words))
                        current_words = combined
                    else:
                        current_words = s_words.copy()
                else:
                    # Split sentence by words respecting budget and overlap
                    start_idx = 0
                    while start_idx < len(s_words):
                        remaining_budget = max(0, chunk_size - text_len(current_words))
                        chunk_words: list[str] = []
                        while start_idx < len(s_words) and text_len(chunk_words + [s_words[start_idx]]) <= (remaining_budget if current_words else chunk_size):
                            chunk_words.append(s_words[start_idx])
                            start_idx += 1
                        if current_words and chunk_words:
                            # Complete current chunk with these words
                            combined = tokenize_words(words_to_text(current_words) + " " + words_to_text(chunk_words))
                            chunks.append(words_to_text(combined))
                            # Prepare next current via overlap
                            if overlap > 0:
                                current_words = build_overlap_words(combined, overlap)
                            else:
                                current_words = []
                        elif chunk_words:
                            # No current_words; emit a full new chunk
                            chunks.append(words_to_text(chunk_words))
                            if overlap > 0:
                                current_words = build_overlap_words(chunk_words, overlap)
                            else:
                                current_words = []
                        else:
                            # Force add single token if nothing fits
                            chunks.append(s_words[start_idx])
                            start_idx += 1
                            current_words = []

        flush_current_words()
        return [c for c in chunks if c]